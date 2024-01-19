import streamlit as st
from streamlit_option_menu import option_menu
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from pptx import Presentation
from features.q_a_pairs import generate_qa_pairs, model, tokenizer

# Set Streamlit page configuration with title, icon, layout, and sidebar state
st.set_page_config(
    page_title="StudyEasy",
    page_icon="https://acmvit.in/logo.png",
    layout="centered",
    initial_sidebar_state="collapsed",
    menu_items=None,
)

# Custom CSS to hide the Streamlit's hamburger menu on the sidebar
st.markdown(
    """
<style>
    [data-testid="collapsedControl"] {
        display: none
    }
</style>
""",
    unsafe_allow_html=True,
)

# Define main function to serve as the entry point of the Streamlit application
def main():
    upload_tab() # Call function to render the upload tab


    st.markdown(
        '<a href="/" target="_self">Back to Home Page</a>', unsafe_allow_html=True
    )


# Function to extract text from a PDF file
def get_pdf_text(pdf):  # Function to read PDF files
    try:
        text = ""

        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()

        text = get_text_chunks(text)
        return text
    except Exception as err:
        from io import StringIO

         # Fallback if PDF text extraction fails, process text as string
        text = " "
        data = pdf.getvalue().decode("utf-8").splitlines()
        for i in range(0, len(data)):
            text += data[i]

        text = get_text_chunks(text)
        return text

# Function to extract text from a PPTX file
def get_pptx_text(pptx_file):  # Function to Read PPTX files
    prs = Presentation(pptx_file)
    text = ""
    for slide_number, slide in enumerate(prs.slides):
        # print(f"Slide {slide_number + 1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # print(shape.text)
                text += shape.text

    return text

# Function to split text into manageable chunks
def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n", chunk_size=1000, chunk_overlap=200, length_function=len
    )
    main_chunks = text_splitter.split_text(text)
    for chunk in range(0, len(main_chunks)):
        main_chunks[chunk] = main_chunks[chunk].replace("\n", "")
    main_chunks = " ".join(main_chunks)
    return main_chunks


def upload_tab():
    st.title("Upload")
    model_selection = st.radio(
        "Choose Functionality",
        ["Summary", "Flashcards"],
        captions=["Generates a summary", "Generates Flashcards"],
    )
    if model_selection == "Summary":
        st.write("You have selected summary.")
    elif model_selection == "Flashcards":
        st.write("You have selected flashcards")

    st.header("Upload Files")
    file = st.file_uploader(
        "Upload PPT or PDF", type=["ppt", "pptx", "pdf"], key="file-uploader"
    )
    submit_file = st.button("Submit File")

    textbox = st.text_input("Enter Content")

    if textbox != "":
        if model_selection == "Flashcards":
            print("Generating flashcards...")
            with st.spinner("Generating Flashcards"):
                flashcards = generate_qa_pairs(textbox, model, tokenizer)
            print("Successfully generated flashcards!!!")
            st.success("Done")
            flashcards_tab(flashcards)
        elif model_selection == "Summary":
            st.write("Summary Model yet to be implemented")

    # You can add functionality to process the uploaded files as needed.
    if file is not None:
        if model_selection == "Flashcards":
            file_type = file.name.split(".")[1]
            if file_type == "pdf":
                st.write("You have uploaded a PDF file.")
                content = get_pdf_text(file)
                print("Received Content")
            elif file_type == "ppt" or file_type == "pptx":
                st.write("You have uploaded a file.")
                content = get_pptx_text(file)
                print("Received Content")

            if submit_file:
                print("Generating flashcards...")
                with st.spinner("Generating Flashcards"):
                    flashcards = generate_qa_pairs(content, model, tokenizer)
                print("Successfully generated flashcards!!!")
                st.success("Done")
                flashcards_tab(flashcards)

        elif model_selection == "Summary":
            st.write("Summary yet to be implemented")


def flashcards_tab(flashcards):
    rotating_card_css = """
    <style>
    .st-emotion-cache-nahz7x {
        perspective: 1000px;
        display: grid;
        grid-template-columns: repeat(3, 1fr);
        grid-gap: 20px;
        justify-content: space-around;
    }

    .card {
        width: 250px; /* Adjusted card width */
        height: 200px;
        margin: 10px;
        border-radius: 15px; /* Added border-radius for rounded edges */
        transform-style: preserve-3d;
        transition: transform 0.8s ease-in-out; /* Adjusted transition properties */
        cursor: pointer;
    }

    .card:hover {
        transform: rotatey(180deg);
        transition: transform 0.8s ease-in-out; /* Added transition property to fix flickering */
    }

    .card .card-inner {
        width: 100%;
        height: 100%;
        position: absolute;
        backface-visibility: hidden;
        display: flex;
        align-items: center;
        justify-content: center;
        text-align: center;
    }

    .card .front {
        background-color: #3498db;
        border-radius: 15px; /* Added border-radius for rounded edges */
    }

    .card .back {
        background-color: #e74c3c;
        transform: rotateY(180deg);
        border-radius: 15px; /* Added border-radius for rounded edges */
    }
    </style>"""
    # html_content = """<div class="card-container">"""
    html_content = ""
    for card in flashcards:
        card_html = f"""
        <div class="card">
            <div class="card-inner front">
            <p>{card['question']}</p>
            </div>
            <div class="card-inner back">
            <p>{card['answer']}</p>
            </div>
        </div>
        """
        html_content += card_html  # Concatenate the card HTML to the main

    st.title("Flash cards")

    st.markdown(rotating_card_css, unsafe_allow_html=True)

    st.markdown(html_content, unsafe_allow_html=True)


if __name__ == "__main__":
    main()
